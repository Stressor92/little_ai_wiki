from __future__ import annotations

from pathlib import Path


def convert_pptx_to_markdown(input_path: Path) -> str:
    """Convert a .pptx presentation into deterministic markdown."""
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ImportError("python-pptx is required: pip install python-pptx") from exc

    prs = Presentation(str(input_path))
    lines: list[str] = []

    title = input_path.stem.replace("_", " ").replace("-", " ").strip()
    lines.append(f"# {title}")

    for idx, slide in enumerate(prs.slides, start=1):
        slide_title = None
        if slide.shapes.title and slide.shapes.title.text:
            slide_title = slide.shapes.title.text.strip()
        if not slide_title:
            slide_title = f"Slide {idx}"

        lines.append("")
        lines.append(f"## Slide {idx}: {slide_title}")

        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            if slide.shapes.title is not None and shape == slide.shapes.title:
                continue

            for paragraph in shape.text_frame.paragraphs:
                text = "".join(run.text for run in paragraph.runs).strip()
                if not text:
                    continue
                level = int(getattr(paragraph, "level", 0) or 0)
                indent = "  " * max(0, level)
                lines.append(f"{indent}- {text}")

        notes_text = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
        if notes_text:
            lines.append("")
            lines.append("### Speaker Notes")
            for note_line in notes_text.splitlines():
                note_line = note_line.strip()
                if note_line:
                    lines.append(f"> {note_line}")

    lines.append("")
    return "\n".join(lines)
